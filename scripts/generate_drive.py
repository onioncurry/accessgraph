# AccessBot — generate the demo "company drive" (real .docx/.pptx/.xlsx/.json files)
# Every file here is referenced by a resource in data/resources.json (file_path),
# so the graph points at artifacts judges can actually open on GitHub.
#   python scripts/generate_drive.py
import json, os, shutil
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DRIVE = os.path.join(ROOT, "company_drive")
os.makedirs(DRIVE, exist_ok=True)

def docx(name, title, paragraphs):
    d = Document()
    d.add_heading(title, level=1)
    for p in paragraphs:
        if isinstance(p, tuple) and p[0] == "h2":
            d.add_heading(p[1], level=2)
        elif isinstance(p, tuple) and p[0] == "bullets":
            for b in p[1]:
                d.add_paragraph(b, style="List Bullet")
        else:
            d.add_paragraph(p)
    d.save(os.path.join(DRIVE, name))

def xlsx(name, sheet, headers, rows, extra_note=None):
    wb = Workbook()
    ws = wb.active
    ws.title = sheet
    ws.append(headers)
    for row in rows:
        ws.append(row)
    if extra_note:
        ws.append([])
        ws.append([extra_note])
    for i, h in enumerate(headers, 1):
        ws.column_dimensions[chr(64 + i)].width = max(14, len(str(h)) + 6)
    wb.save(os.path.join(DRIVE, name))

def pptx(path, slides):
    prs = Presentation()
    for i, (title, bullets) in enumerate(slides):
        layout = prs.slide_layouts[0 if i == 0 else 1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        if i == 0 and len(s.placeholders) > 1:
            s.placeholders[1].text = bullets[0] if bullets else ""
        elif bullets:
            body = s.placeholders[1].text_frame
            body.text = bullets[0]
            for b in bullets[1:]:
                para = body.add_paragraph()
                para.text = b
                para.level = 0
    prs.save(path)

# ---------------- Phoenix files (attached to existing resources) ----------------
docx("Project Phoenix - Product Brief.docx", "Project Phoenix — Product Brief", [
    "Owner: Shota Gushima · Status: In flight · Confidentiality: Internal",
    ("h2", "Problem"),
    "New users drop out during onboarding: 38% never reach the first success moment. Activation is the #1 blocker to Q2 revenue goals.",
    ("h2", "Goal"),
    "Ship Onboarding v2 and cut step-3 drop-off from 38% to under 15% by the end of Q2.",
    ("h2", "Success criteria"),
    ("bullets", ["Activation rate ≥ 55% (from 41%)", "Time-to-first-value under 4 minutes", "Support tickets tagged 'onboarding' down 50%"]),
    ("h2", "Scope"),
    ("bullets", ["Redesigned flow (see Onboarding Flow v2 in Figma)", "Fix PROD-1327 (step-3 crash)", "New progress milestones tracked in Q2 Roadmap"]),
])
xlsx("Q2 Roadmap & Milestones.xlsx", "Q2 Roadmap",
    ["Milestone", "Owner", "Due", "Status", "Notes"],
    [["Onboarding v2 design freeze", "Koichi Ikeno", "2026-04-18", "Done", "Figma v2 approved"],
     ["PROD-1327 fix shipped", "Mitsuhiro Suzuki", "2026-05-09", "In progress", "Blocking launch"],
     ["Beta cohort (50 users)", "Shota Gushima", "2026-05-30", "In progress", "Recruiting via #product-dev"],
     ["Activation dashboard live", "Rei Kawaji", "2026-06-06", "Not started", "Needs analytics access"],
     ["GA launch", "Shota Gushima", "2026-06-27", "Not started", "Gate: beta NPS ≥ 40"]])
docx("Onboarding Flow v2 - Design Notes.docx", "Onboarding Flow v2 — Design Notes", [
    "Owner: Koichi Ikeno · Source of truth: Figma (this doc mirrors decisions)",
    ("h2", "Key decisions"),
    ("bullets", ["3 steps instead of 7 — progressive profiling moves to week 2",
                 "Template gallery on step 2 (top 6 by usage)",
                 "Skip is always visible; empty state teaches the core action"]),
    ("h2", "Open questions"),
    ("bullets", ["Does SSO-first reduce step-1 abandonment?", "Mobile: native sheet vs web view"]),
])
docx("PROD-1327 - Bug Report.docx", "PROD-1327 — Fix onboarding flow issue", [
    "Reporter: Shota Gushima · Assignee: Mitsuhiro Suzuki · Priority: P1 · Status: In progress",
    ("h2", "Symptom"),
    "Users drop at step 3 of onboarding. Client throws TypeError when the workspace list is empty; the continue button never enables.",
    ("h2", "Repro"),
    ("bullets", ["Fresh account, no invites", "Complete steps 1-2", "Step 3 renders spinner forever (console: cannot read length of undefined)"]),
    ("h2", "Impact"),
    "Affects 100% of solo signups — the exact cohort Onboarding v2 targets. Blocking the Phoenix launch gate.",
])

# ---------------- Sales / GTM (real Crustdata data goes here) ----------------
with open(os.path.join(ROOT, "data", "crustdata_companies.json"), encoding="utf-8-sig") as f:
    companies = json.load(f)
with open(os.path.join(ROOT, "data", "crustdata_demo.json"), encoding="utf-8") as f:
    stripe = json.load(f)["company"]
rows = [[stripe["name"], "stripe.com", stripe["headquarters"], stripe["employee_count_range"], str(stripe["year_founded"])[:4], "Customer", "$120,000", "Renewal Q3"]]
for c in companies:
    if c.get("name") and c["name"] != "FETCH_FAILED":
        rows.append([c["name"], c["domain"], c.get("hq") or "—", c.get("employees") or "—", c.get("founded") or "—", "Prospect", "—", "Discovery call scheduled"])
xlsx("Customer Accounts.xlsx", "Accounts",
    ["Company", "Domain", "HQ", "Employees", "Founded", "Stage", "ARR", "Next step"],
    rows,
    "Company firmographics pulled LIVE from the Crustdata API (see data/crustdata_companies.json). Stage/ARR columns are fictional demo values.")
pptx(os.path.join(DRIVE, "Phoenix Launch Deck.pptx"), [
    ("Phoenix Launch", ["GTM narrative & announcement plan — Internal"]),
    ("The story", ["Onboarding used to take 7 steps. Now it takes 3.",
                    "Activation +14pts in beta", "Launch moment: changelog + founder thread + lifecycle email"]),
    ("Rollout", ["Week 1: beta cohort GA", "Week 2: all new signups", "Week 3: existing workspaces prompt"]),
    ("Asks", ["Sales: update demo script", "Support: macro refresh", "Everyone: amplify launch thread"]),
])
pptx(os.path.join(DRIVE, "Q2 Sales Pipeline Review.pptx"), [
    ("Q2 Pipeline Review", ["Sales / GTM — Internal"]),
    ("Headlines", ["Coverage 3.1x (target 3x)", "Win rate 24% (+3pts QoQ)", "Slippage risk: 2 enterprise deals waiting on security review"]),
    ("Focus", ["Mid-market expansion plays", "Onboarding v2 as a wedge story", "Tighten stage-2 exit criteria"]),
])
shutil.copyfile(os.path.join(ROOT, "data", "crustdata_companies.json"),
                os.path.join(DRIVE, "Market Research - Competitor Scan.json"))

# ---------------- People Ops / Finance (confidential in the graph) ----------------
xlsx("2026 Hiring Plan.xlsx", "Headcount",
    ["Team", "Q1", "Q2", "Q3", "Q4", "Notes"],
    [["Engineering", 2, 3, 2, 2, "Backend-heavy H1"],
     ["Design", 1, 1, 0, 1, "Design Lead search active"],
     ["Product", 0, 1, 1, 0, "PM after GA"],
     ["GTM", 1, 2, 2, 3, "Scale post-launch"]])
docx("Interview Notes - Design Lead Candidate.docx", "Interview Notes — Design Lead Candidate (CONFIDENTIAL)", [
    "HR data — access via AccessBot only. Approver: hiring manager.",
    ("h2", "Panel summary"),
    ("bullets", ["Portfolio: strong systems thinking, weaker motion", "Craft exercise: 8/10", "Values interview: clear hire signal"]),
    ("h2", "Decision"),
    "Proceed to offer discussion. Comp band per 2026 Hiring Plan.",
])
xlsx("Q2 Payroll.xlsx", "Payroll (demo)",
    ["Employee", "Role", "Monthly (demo)", "Notes"],
    [["<fictional demo data>", "—", "—", "Payroll numbers are intentionally fictional"],
     ["Shota Gushima", "Product Owner", "¥—", "demo placeholder"],
     ["Rei Kawaji", "Product Associate", "¥—", "demo placeholder"]],
    "CONFIDENTIAL in the access graph (hr_data policy: manager approval, 3-day grants, never auto-bundled).")
pptx(os.path.join(DRIVE, "Q2 Board Update.pptx"), [
    ("Q2 Board Update", ["CONFIDENTIAL — finance_data policy"]),
    ("Metrics", ["ARR growth steady", "Activation is the constraint — Phoenix addresses it", "Runway: healthy"]),
    ("Asks", ["Intro to design-systems advisor", "Enterprise security review fast-track"]),
])

# ---------------- Brand / Operations ----------------
docx("Brand Guidelines v3.docx", "Brand Guidelines v3", [
    "Owner: Koichi Ikeno · Applies to all external surfaces",
    ("h2", "Logo"), ("bullets", ["Clear space = 1x mark height", "Never recolor the mark", "Dark backgrounds use the mono variant"]),
    ("h2", "Type & color"), ("bullets", ["Display: Inter Tight", "Body: Inter", "Primary #4A154B, success #007A5A"]),
])
docx("Company Handbook.docx", "Company Handbook", [
    "Owner: Operations · Shared with everyone",
    ("h2", "How we work"), ("bullets", ["Async-first; decisions in writing", "Meetings have an owner and a doc", "Access requests go through AccessBot — least privilege by default"]),
    ("h2", "Time off"), ("bullets", ["Unlimited PTO with a 15-day minimum", "Company recharge week in August"]),
])
docx("Incident Response Runbook.docx", "Incident Response Runbook", [
    "Owner: Mitsuhiro Suzuki · On-call reference",
    ("h2", "Sev1 first 15 minutes"),
    ("bullets", ["Page secondary + comms lead", "Open #incident channel and status doc", "Mitigate first, diagnose second", "Customer comms at 30-minute cadence"]),
    ("h2", "After"), ("bullets", ["Blameless postmortem within 5 days", "Action items tracked in Jira with owners"]),
])

# ---------------- Judge-facing pitch deck (docs/) ----------------
pptx(os.path.join(ROOT, "docs", "AccessBot - Pitch Deck.pptx"), [
    ("AccessBot", ["A G-Brain skill that turns a Slack task into a minimum access package — for humans and AI agents."]),
    ("Problem", ["Work is assigned in Slack. Access lives in ten tools.",
                  "Every handoff starts with: I cannot open the roadmap. Who owns the Jira board?",
                  "AI agents cannot even ask — no machine-readable permissions = no safe automation."]),
    ("Solution", ["Mention @AccessBot on a task", "It parses assignee, intent, and project",
                   "Queries the G-Brain access graph (people, files, policies, past tasks)",
                   "Returns the least-privilege package: what, level, duration, approver — one click to grant"]),
    ("Why it is safe", ["Per-resource least privilege (write only on work surfaces)",
                          "Confidential data never auto-bundled — blocked unless justified",
                          "Fail-closed, deterministic, injection-inert", "120+ adversarial tests, score 100/100"]),
    ("Why RFS — all three", ["Company Brain: knowledge, policies, AND the skill live in the brain",
                               "Software for Agents: strict JSON before an agent acts",
                               "Dynamic Software Interfaces: per-task access UI in the flow of work"]),
    ("It runs on real data", ["People layer generated from the Crustdata API (real Stripe org, 679 matches)",
                                "Same unchanged skill resolves access for real employees",
                                "Fully offline demo: all data in-repo, zero API dependency"]),
    ("Ask", ["Try it: onioncurry.github.io/accessgraph", "Repo: github.com/onioncurry/accessgraph"]),
])

print("company_drive/ written:", len(os.listdir(DRIVE)), "files + docs/AccessBot - Pitch Deck.pptx")
